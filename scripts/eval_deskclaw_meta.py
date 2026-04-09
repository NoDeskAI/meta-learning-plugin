#!/usr/bin/env python3
"""DeskClaw Meta-Learning Automated Evaluation.

10 scenarios × 5 task domains = 50 cases.
Drives a real DeskClaw instance via Gateway HTTP API to verify
that meta-learning preferences transfer across task domains.

Usage:
    python scripts/eval_deskclaw_meta.py
    python scripts/eval_deskclaw_meta.py --scenarios backup_first,output_dir
    python scripts/eval_deskclaw_meta.py --domains shell,excel
    python scripts/eval_deskclaw_meta.py --skip-teach
    python scripts/eval_deskclaw_meta.py --gateway http://127.0.0.1:18790

Requires:
    - DeskClaw instance running (Gateway on port 18790)
    - NODESK_API_KEY env var (for LLM Judge)
    - Python: httpx, openpyxl, python-docx, python-pptx
"""

from __future__ import annotations

import os as _os

for _pv in (
    "http_proxy", "https_proxy", "all_proxy",
    "HTTP_PROXY", "HTTPS_PROXY", "ALL_PROXY",
):
    _os.environ.pop(_pv, None)

import argparse
import asyncio
import json
import logging
import os
import shutil
import struct
import time
import uuid
import zlib
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

import httpx

import sys

SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent
sys.path.insert(0, str(PROJECT_ROOT / "src"))

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(name)s: %(message)s",
)
logger = logging.getLogger("eval_deskclaw_meta")

RESULTS_BASE = PROJECT_ROOT / "abtest" / "results" / "deskclaw_meta"

DEFAULT_GATEWAY = "http://127.0.0.1:18790"
META_DATA_DIR = Path("~/.deskclaw/nanobot/workspace/meta-learning-data").expanduser()
NANOBOT_SKILLS_DIR = Path("~/.deskclaw/nanobot/workspace/skills").expanduser()
DOCUMENTS_DIR = Path("~/Documents").expanduser()

NODESK_GATEWAY = "https://llm-gateway-api.nodesk.tech/default/v1"
JUDGE_MODEL = "anthropic/claude-opus-4.6"

ALL_DOMAINS = ("shell", "coding", "excel", "word", "pptx")


# =====================================================================
# Gateway Client
# =====================================================================

class DeskClawClient:
    """Async HTTP client for the DeskClaw Gateway."""

    def __init__(self, base_url: str = DEFAULT_GATEWAY, timeout: float = 600):
        self.base_url = base_url.rstrip("/")
        self._client = httpx.AsyncClient(timeout=timeout)

    async def health_check(self) -> bool:
        try:
            resp = await self._client.get(f"{self.base_url}/health")
            return resp.status_code == 200
        except Exception:
            return False

    async def send_message(self, message: str, session_id: str) -> dict[str, Any]:
        """POST /chat — blocks until the agent finishes responding."""
        resp = await self._client.post(
            f"{self.base_url}/chat",
            json={"message": message, "session_id": session_id},
        )
        resp.raise_for_status()
        return resp.json()

    async def get_history(self, session_id: str) -> list[dict[str, Any]]:
        resp = await self._client.get(
            f"{self.base_url}/chat/{session_id}/history",
        )
        resp.raise_for_status()
        return resp.json()

    async def delete_session(self, session_id: str) -> None:
        try:
            await self._client.delete(f"{self.base_url}/chat/{session_id}")
        except Exception:
            pass

    async def wait_learning_done(
        self, timeout: float = 180, after_ts: float | None = None,
    ) -> bool:
        """Poll layer2_state.json and SKILL.md until learning is fully applied.

        Returns True only when layer2 status == completed AND the nanobot
        SKILL.md has been updated (mtime > after_ts), meaning sync_nobot
        ran and the agent will see the new rules in subsequent sessions.
        """
        state_path = META_DATA_DIR / "signal_buffer" / "layer2_state.json"
        skill_path = NANOBOT_SKILLS_DIR / "meta-learning" / "SKILL.md"
        deadline = time.time() + timeout
        layer2_done = False
        while time.time() < deadline:
            if not layer2_done and state_path.exists():
                if after_ts and state_path.stat().st_mtime < after_ts:
                    await asyncio.sleep(3)
                    continue
                try:
                    data = json.loads(state_path.read_text())
                    status = data.get("status", "")
                    if status == "completed":
                        logger.info("    layer2 pipeline completed")
                        layer2_done = True
                    elif status == "running":
                        logger.debug("    layer2 still running …")
                except (json.JSONDecodeError, OSError):
                    pass

            if layer2_done:
                if skill_path.exists() and (
                    not after_ts or skill_path.stat().st_mtime > after_ts
                ):
                    logger.info("    SKILL.md updated — learning fully applied")
                    return True
                logger.debug("    layer2 done, waiting for SKILL.md sync …")

            await asyncio.sleep(3)
        if layer2_done:
            logger.warning("    layer2 completed but SKILL.md not updated within %.0fs", timeout)
        else:
            logger.warning("    layer2 pipeline timed out after %.0fs", timeout)
        return layer2_done

    async def close(self):
        await self._client.aclose()


# =====================================================================
# Seed File Generators
# =====================================================================

def _prepare_seed_files():
    """Create all seed files needed by evaluation scenarios."""
    doc = DOCUMENTS_DIR
    doc.mkdir(parents=True, exist_ok=True)

    _create_sample_xlsx(doc / "Q3销售报表.xlsx")
    _create_budget_xlsx(doc / "budget.xlsx")
    _create_sample_csv(doc / "sales_data.csv")
    _create_data_csv(doc / "data.csv")

    _create_sample_docx(doc / "合作协议.docx")
    _create_contract_docx(doc / "合同.docx")
    _create_product_docx(doc / "产品说明.docx")

    _create_sample_pptx(doc / "公司介绍.pptx")
    _create_deck_pptx(doc / "company-deck.pptx")
    _create_training_pptx(doc / "培训资料.pptx")

    _create_report_json(doc / "report.json")
    _create_readme_md(doc / "README.md")
    _create_sample_png(doc / "photo.png")
    _create_main_py(doc / "main.py")
    _create_utils_js(doc / "utils.js")

    _create_broken_app(doc / "app.py")
    _create_corrupted_xlsx(doc / "corrupted.xlsx")
    _create_old_doc(doc / "old_format.doc")
    _create_broken_pptx(doc / "broken.pptx")

    drafts = doc / "drafts"
    drafts.mkdir(exist_ok=True)
    for name in ("草稿1.docx", "草稿2.docx", "草稿3.docx"):
        _create_simple_docx(drafts / name)

    old_proj = Path.home() / "old-project"
    old_proj.mkdir(exist_ok=True)
    (old_proj / "README.md").write_text("# Legacy Project\n")
    (old_proj / "main.py").write_text("print('hello')\n")
    (old_proj / "data.csv").write_text("a,b\n1,2\n")

    logger.info("Seed files prepared in %s", doc)


def _create_sample_xlsx(path: Path):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Q3 Sales"
    ws.append(["Month", "Product", "Revenue", "Quantity"])
    ws.append(["July", "Widget A", 15000, 300])
    ws.append(["August", "Widget A", 18000, 360])
    ws.append(["September", "Widget A", "=SUM(C2:C3)", "=SUM(D2:D3)"])
    wb.save(path)


def _create_budget_xlsx(path: Path):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Category", "Q1", "Q2", "Q3", "Q4"])
    ws.append(["Marketing", 50000, 55000, 60000, 65000])
    ws.append(["Engineering", 200000, 210000, 220000, 230000])
    ws.append(["Operations", 80000, 82000, 84000, 86000])
    wb.save(path)


def _create_sample_csv(path: Path):
    lines = [
        "date,product,region,amount,quantity",
        "2026-01-15,Widget A,North,1200,24",
        "2026-01-20,Widget B,South,800,16",
        "2026-02-10,Widget A,East,1500,30",
        "2026-02-28,Widget C,West,2200,44",
        "2026-03-05,Widget B,North,950,19",
        "2026-03-18,Widget A,South,1800,36",
    ]
    path.write_text("\n".join(lines) + "\n")


def _create_data_csv(path: Path):
    lines = [
        "id,name,score,grade",
        "1,Alice,92,A",
        "2,Bob,78,B",
        "3,Charlie,85,B+",
        "4,Diana,96,A+",
        "5,Eve,70,C+",
    ]
    path.write_text("\n".join(lines) + "\n")


def _create_sample_docx(path: Path):
    from docx import Document
    doc = Document()
    doc.add_heading("合作协议", level=1)
    doc.add_paragraph("甲方：ABC 公司")
    doc.add_paragraph("乙方：XYZ 公司")
    doc.add_paragraph("付款周期：30 天")
    doc.add_paragraph("合作期限：2026年1月1日 至 2026年12月31日")
    doc.save(str(path))


def _create_contract_docx(path: Path):
    from docx import Document
    doc = Document()
    doc.add_heading("服务合同", level=1)
    doc.add_paragraph("第一条：服务范围")
    doc.add_paragraph("甲方委托乙方提供技术咨询服务。")
    doc.add_paragraph("第二条：费用及支付")
    doc.add_paragraph("服务费用总计人民币叁拾万元整（¥300,000）。")
    doc.add_paragraph("第三条：保密条款")
    doc.add_paragraph("双方对合作过程中获知的商业秘密负有保密义务。")
    doc.save(str(path))


def _create_product_docx(path: Path):
    from docx import Document
    doc = Document()
    doc.add_heading("产品说明书", level=1)
    doc.add_paragraph("产品名称：SmartWidget Pro")
    doc.add_paragraph("版本：2.0")
    doc.add_paragraph("功能特性：")
    for feat in ["自动化数据处理", "多平台兼容", "实时协作", "安全加密"]:
        doc.add_paragraph(feat, style="List Bullet")
    doc.save(str(path))


def _create_simple_docx(path: Path):
    from docx import Document
    doc = Document()
    doc.add_paragraph(f"Draft: {path.stem}")
    doc.add_paragraph("This is a draft document for testing.")
    doc.save(str(path))


def _create_sample_pptx(path: Path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "公司介绍"
    slide.placeholders[1].text = "ABC 科技有限公司"
    for title in ("团队介绍", "组织架构", "核心业务", "联系方式"):
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = title
        sl.placeholders[1].text = f"{title}内容占位"
    prs.save(str(path))


def _create_deck_pptx(path: Path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Company Overview"
    slide.placeholders[1].text = "FY2026 Strategy Deck"
    for title in ("Mission", "Products", "Roadmap"):
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = title
        sl.placeholders[1].text = f"{title} placeholder content"
    prs.save(str(path))


def _create_training_pptx(path: Path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "培训资料"
    slide.placeholders[1].text = "2026 年度技术培训"
    for title in ("培训目标", "课程大纲", "实操环节", "考核标准"):
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = title
        sl.placeholders[1].text = f"{title}详细内容"
    prs.save(str(path))


def _create_report_json(path: Path):
    data = {
        "report_date": "2026-03-31",
        "department": "Engineering",
        "metrics": [
            {"name": "commits", "value": 342},
            {"name": "pull_requests", "value": 87},
            {"name": "bugs_fixed", "value": 56},
            {"name": "features_shipped", "value": 12},
        ],
    }
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False))


def _create_readme_md(path: Path):
    path.write_text(
        "# Sample Project\n\n"
        "## Quick Start\n\n"
        "```bash\npip install -r requirements.txt\npython main.py\n```\n\n"
        "## Features\n\n- Feature A\n- Feature B\n- Feature C\n"
    )


def _create_sample_png(path: Path):
    """Create a minimal 8×8 red PNG without Pillow."""
    width, height = 8, 8
    raw_rows = []
    for _ in range(height):
        row = b"\x00"
        for _ in range(width):
            row += b"\xff\x00\x00"
        raw_rows.append(row)
    raw_data = b"".join(raw_rows)

    def _chunk(ctype: bytes, data: bytes) -> bytes:
        c = ctype + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    png = b"\x89PNG\r\n\x1a\n"
    png += _chunk(b"IHDR", ihdr)
    png += _chunk(b"IDAT", zlib.compress(raw_data))
    png += _chunk(b"IEND", b"")
    path.write_bytes(png)


def _create_main_py(path: Path):
    path.write_text(
        "import sqlite3\n\n"
        "DB_PATH = 'app.db'\n\n"
        "def get_users():\n"
        "    conn = sqlite3.connect(DB_PATH)\n"
        "    cursor = conn.cursor()\n"
        "    cursor.execute('SELECT * FROM users')\n"
        "    rows = cursor.fetchall()\n"
        "    conn.close()\n"
        "    return rows\n\n"
        "def add_user(name, email):\n"
        "    conn = sqlite3.connect(DB_PATH)\n"
        "    cursor = conn.cursor()\n"
        "    cursor.execute('INSERT INTO users (name, email) VALUES (?, ?)', (name, email))\n"
        "    conn.commit()\n"
        "    conn.close()\n\n"
        "if __name__ == '__main__':\n"
        "    print(get_users())\n"
    )


def _create_utils_js(path: Path):
    path.write_text(
        "function formatDate(date) {\n"
        "  const d = new Date(date);\n"
        "  return d.toISOString().split('T')[0];\n"
        "}\n\n"
        "function capitalize(str) {\n"
        "  return str.charAt(0).toUpperCase() + str.slice(1);\n"
        "}\n\n"
        "function debounce(fn, ms) {\n"
        "  let timer;\n"
        "  return function(...args) {\n"
        "    clearTimeout(timer);\n"
        "    timer = setTimeout(() => fn.apply(this, args), ms);\n"
        "  };\n"
        "}\n\n"
        "module.exports = { formatDate, capitalize, debounce };\n"
    )


def _create_broken_app(path: Path):
    path.write_text(
        "from flask import Flask\n"
        "from database_manager import DatabaseManager\n\n"
        "app = Flask(__name__)\n"
        "db = DatabaseManager('production.db')\n\n"
        "@app.route('/')\n"
        "def index():\n"
        "    return db.get_status()\n\n"
        "if __name__ == '__main__':\n"
        "    app.run(port=5000)\n"
    )


def _create_corrupted_xlsx(path: Path):
    path.write_bytes(b"PK\x03\x04" + b"\x00" * 100 + b"corrupted xlsx content")


def _create_old_doc(path: Path):
    header = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
    path.write_bytes(header + b"\x00" * 200 + b"Old format document content")


def _create_broken_pptx(path: Path):
    path.write_bytes(b"PK\x03\x04" + b"\x00" * 80 + b"broken pptx data")


# =====================================================================
# Data Models
# =====================================================================

@dataclass
class TreatmentCase:
    domain: str
    prompt: str


@dataclass
class Scenario:
    name: str
    preference: str
    teach_domain: str
    teach_task: str
    correction: str
    treatment_cases: list[TreatmentCase]
    verify_keywords: list[str] = field(default_factory=list)


# =====================================================================
# 10 Scenarios × 5 Domains
# =====================================================================

SCENARIOS: list[Scenario] = [
    # ---- 1. 备份优先 ----
    Scenario(
        name="backup_first",
        preference="修改重要文件/文档前先备份原件",
        teach_domain="shell",
        teach_task="帮我修改 ~/.bashrc，添加一个 alias：`alias ll='ls -la'`",
        correction=(
            "你直接改了文件。以后修改任何重要文件前，先用 cp 备份一份 .bak，"
            "比如 cp ~/.bashrc ~/.bashrc.bak，再修改。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我修改 /etc/hosts，添加一条 `127.0.0.1 myapp.local` 的记录"),
            TreatmentCase("coding", "帮我重构 ~/Documents/main.py，把数据库连接逻辑从内联改成单独的 db.py 模块"),
            TreatmentCase("excel",  "帮我更新 ~/Documents/Q3销售报表.xlsx，把9月数据行的公式从 SUM 改成 SUMIF"),
            TreatmentCase("word",   "帮我编辑 ~/Documents/合作协议.docx，把付款周期从 30 天改成 45 天"),
            TreatmentCase("pptx",   "帮我修改 ~/Documents/公司介绍.pptx，更新第3页的组织架构图"),
        ],
        verify_keywords=["备份", "backup", "bak", "cp "],
    ),
    # ---- 2. 输出目录 ----
    Scenario(
        name="output_dir",
        preference="所有新生成的文件统一放到 ~/Documents/output/ 目录下",
        teach_domain="coding",
        teach_task="帮我创建一个 Python 项目 calculator，包含 main.py 和 README.md",
        correction=(
            "你把项目创建在了当前目录。我的习惯是所有新生成的文件和项目统一放到"
            " ~/Documents/output/ 下面，请以后都放那里。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我写一个 backup.sh 脚本，用来备份指定目录"),
            TreatmentCase("coding", "帮我创建一个 Node.js todo-app 项目，包含 index.js 和 package.json"),
            TreatmentCase("excel",  "帮我做一个本月团队工时统计表"),
            TreatmentCase("word",   "帮我写一份《Q2 项目进展报告》"),
            TreatmentCase("pptx",   "帮我做一个《新产品发布会》的 PPT，5 页左右"),
        ],
        verify_keywords=["output", "~/Documents/output"],
    ),
    # ---- 3. 操作确认 ----
    Scenario(
        name="confirm_before_act",
        preference="执行不可逆操作前先列出影响范围并等待确认",
        teach_domain="shell",
        teach_task="帮我删除 ~/old-project 文件夹",
        correction=(
            "你直接 rm -rf 了！以后执行删除、覆盖这种不可逆操作之前，"
            "先 ls 列出内容告诉我有什么，等我确认后再执行。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我清理 /tmp 下超过 7 天的临时文件"),
            TreatmentCase("coding", "帮我用 git reset --hard HEAD~3 回退最近 3 个 commit"),
            TreatmentCase("excel",  "帮我用新模板替换现有的 ~/Documents/budget.xlsx"),
            TreatmentCase("word",   "帮我删除 ~/Documents/drafts/ 目录下的所有文档草稿"),
            TreatmentCase("pptx",   "帮我用新品牌模板替换现有的 ~/Documents/company-deck.pptx"),
        ],
        verify_keywords=["确认", "列出", "confirm", "以下"],
    ),
    # ---- 4. 中文输出 ----
    Scenario(
        name="chinese_output",
        preference="生成内容（标题、表头、文案、注释）默认使用中文",
        teach_domain="excel",
        teach_task="帮我做一个员工信息登记表，包含姓名、部门、入职日期",
        correction=(
            "表头和 Sheet 名称你用了英文。以后所有生成内容请默认使用中文，"
            "包括标题、表头、Sheet名、文件名等。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我写一个 cron 定时任务配置，每天凌晨 2 点清理日志，加上注释说明"),
            TreatmentCase("coding", "帮我写一个 Python 日志轮转脚本，加上详细的代码注释"),
            TreatmentCase("excel",  "帮我做一个项目进度跟踪表，包含任务名、负责人、状态、截止日期"),
            TreatmentCase("word",   "帮我写一份《技术部周会纪要》模板"),
            TreatmentCase("pptx",   "帮我做一个《新员工入职培训》PPT，6 页"),
        ],
        verify_keywords=["中文"],
    ),
    # ---- 5. 完成验证 ----
    Scenario(
        name="verify_result",
        preference="操作完成后必须验证结果正确性",
        teach_domain="coding",
        teach_task="帮我安装 requests 库",
        correction=(
            "你安装完就结束了，没有验证。以后操作完成后必须验证结果，"
            "比如安装包后用 pip show 确认、生成文件后检查文件是否存在并看看内容是否正确。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我安装 jq 命令行工具"),
            TreatmentCase("coding", "帮我用 pip 安装 pyyaml 库"),
            TreatmentCase("excel",  "帮我生成一个带 SUM 公式的部门费用汇总表"),
            TreatmentCase("word",   "帮我生成一份 10 页的《年度技术总结报告》"),
            TreatmentCase("pptx",   "帮我生成一个 5 页的《产品功能介绍》PPT"),
        ],
        verify_keywords=["验证", "verify", "确认", "检查", "pip show", "ls "],
    ),
    # ---- 6. 简洁回复 ----
    Scenario(
        name="concise_reply",
        preference="对执行结果做简洁总结，不堆砌原始输出",
        teach_domain="shell",
        teach_task="帮我查看系统所有 Python 相关的版本信息",
        correction=(
            "你把命令输出全贴出来了，太长了。以后不要把原始输出直接甩给我，"
            "提取关键信息做简洁的总结就好。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我查看系统磁盘使用情况"),
            TreatmentCase("coding", "帮我运行 pytest 测试套件，告诉我结果"),
            TreatmentCase("excel",  "帮我分析 ~/Documents/sales_data.csv 的基本统计信息（总额、均值、最大最小）"),
            TreatmentCase("word",   "帮我提取 ~/Documents/合同.docx 的关键条款摘要"),
            TreatmentCase("pptx",   "帮我分析 ~/Documents/培训资料.pptx 的内容结构概要"),
        ],
        verify_keywords=["总结", "简洁", "summary"],
    ),
    # ---- 7. 文件命名 ----
    Scenario(
        name="file_naming",
        preference="输出文件名包含日期和类型标识，格式 YYYYMMDD_类型_名称",
        teach_domain="word",
        teach_task="帮我写一份关于上周项目进展的周报",
        correction=(
            "文件名太随意了。以后生成的文件名请遵循格式：YYYYMMDD_类型_名称，"
            "比如 20260403_周报_项目进展.docx。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我备份 nginx 配置文件到 ~/backup/"),
            TreatmentCase("coding", "帮我生成一份 API 接口文档，保存为 Markdown"),
            TreatmentCase("excel",  "帮我做一份本月开支记录表"),
            TreatmentCase("word",   "帮我写一份《技术调研：向量数据库选型》的报告"),
            TreatmentCase("pptx",   "帮我做一个《Q2 项目启动会》PPT"),
        ],
        verify_keywords=["2026", "_"],
    ),
    # ---- 8. 分步执行 ----
    Scenario(
        name="step_by_step",
        preference="复杂任务先列出步骤计划，等确认后再执行",
        teach_domain="pptx",
        teach_task=(
            "帮我做一个完整的产品路线图 PPT，包含现状分析、目标规划、"
            "里程碑时间线、资源需求、风险评估"
        ),
        correction=(
            "你一口气全做完了，但有些内容不是我想要的。以后遇到复杂任务，"
            "先列出你的步骤计划和每一步要做的内容，等我确认后再开始执行。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我搭建一个完整的 Docker + Nginx + Python 开发环境"),
            TreatmentCase("coding", "帮我搭建一个完整的 Flask REST API 项目骨架，包含路由、模型、测试"),
            TreatmentCase("excel",  "帮我做一个完整的年度财务预算模型，包含收入、成本、利润预测"),
            TreatmentCase("word",   "帮我写一份完整的《商业计划书》，包含市场分析、产品方案、财务预测"),
            TreatmentCase("pptx",   "帮我做一个 20 页的《年度工作总结》PPT，包含业绩回顾、项目亮点、团队建设、明年计划"),
        ],
        verify_keywords=["步骤", "计划", "大纲", "plan", "step"],
    ),
    # ---- 9. 保留原始 ----
    Scenario(
        name="preserve_original",
        preference="做格式转换时保留源文件，不删除或覆盖",
        teach_domain="excel",
        teach_task="帮我把 ~/Documents/data.csv 转成 Excel 格式",
        correction=(
            "你把原始 CSV 删掉了！做格式转换时要保留原始文件，"
            "只新建目标格式的文件，不要删除或覆盖源文件。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我把 ~/Documents/photo.png 转成 JPG 格式"),
            TreatmentCase("coding", "帮我把 ~/Documents/utils.js 转成 TypeScript（utils.ts）"),
            TreatmentCase("excel",  "帮我把 ~/Documents/report.json 的数据导入到 Excel 表格中"),
            TreatmentCase("word",   "帮我把 ~/Documents/README.md 转换成 Word 文档"),
            TreatmentCase("pptx",   "帮我把 ~/Documents/产品说明.docx 的内容转成 PPT 演示稿"),
        ],
        verify_keywords=["保留", "原始", "不删除", "preserve"],
    ),
    # ---- 10. 错误处理 ----
    Scenario(
        name="error_handling",
        preference="遇到错误给出清晰的诊断和修复建议，不静默忽略",
        teach_domain="coding",
        teach_task="帮我运行 python3 ~/Documents/app.py 启动服务",
        correction=(
            "报错了你只说了一句'启动失败'就结束了。遇到错误时请给出：1)具体错误原因 "
            "2)可能的影响 3)修复建议，不要一句话带过。"
        ),
        treatment_cases=[
            TreatmentCase("shell",  "帮我启动 redis-server 服务"),
            TreatmentCase("coding", "帮我运行 npm install 然后 npm start 启动前端项目"),
            TreatmentCase("excel",  "帮我打开并分析 ~/Documents/corrupted.xlsx 的数据"),
            TreatmentCase("word",   "帮我转换 ~/Documents/old_format.doc 为 docx 格式"),
            TreatmentCase("pptx",   "帮我读取 ~/Documents/broken.pptx 的内容"),
        ],
        verify_keywords=["原因", "修复", "建议", "error", "fix"],
    ),
]

SCENARIO_MAP: dict[str, Scenario] = {s.name: s for s in SCENARIOS}


# =====================================================================
# LLM Judge
# =====================================================================

async def llm_judge(
    preference: str,
    response_text: str,
    api_key: str,
) -> dict[str, Any]:
    """Ask an LLM whether *response_text* complies with *preference*."""
    prompt = (
        "你是一个行为合规评估专家。判断以下 AI Agent 的回复是否遵循了用户的偏好规则。\n\n"
        f"【用户偏好】\n{preference}\n\n"
        f"【Agent 回复全文】\n{response_text[:6000]}\n\n"
        "请判断 Agent 的行为是否符合用户偏好。只返回 JSON，不要 markdown 代码块：\n"
        '{"compliant": true/false, "confidence": 0.0-1.0, "reasoning": "简要说明"}'
    )

    try:
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                f"{NODESK_GATEWAY}/chat/completions",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": JUDGE_MODEL,
                    "messages": [
                        {
                            "role": "system",
                            "content": "你是行为合规评估专家。只返回合法 JSON，不加 markdown。",
                        },
                        {"role": "user", "content": prompt},
                    ],
                    "temperature": 0,
                },
            )
            resp.raise_for_status()
            raw = resp.json()["choices"][0]["message"]["content"].strip()
            if raw.startswith("```"):
                lines = raw.split("\n")
                raw = "\n".join(lines[1:-1])
            data = json.loads(raw)
            return {
                "compliant": bool(data.get("compliant", False)),
                "confidence": float(data.get("confidence", 0.5)),
                "reasoning": data.get("reasoning", ""),
            }
    except Exception as e:
        logger.warning("LLM Judge failed: %s", e)
        return {"compliant": None, "confidence": 0.0, "reasoning": f"judge error: {e}"}


# =====================================================================
# Pipeline Artifact Verification (Teaching Phase)
# =====================================================================

def _skill_has_learned_rules(content: str) -> bool:
    """True if SKILL.md contains actual learned rules (not just bootstrap)."""
    in_rules_section = False
    for line in content.splitlines():
        if line.startswith("# Meta-Learning Rules"):
            in_rules_section = True
            continue
        if in_rules_section:
            if line.startswith("#"):
                break
            if line.startswith("- "):
                return True
    return False


def check_pipeline_artifacts(verify_keywords: list[str] | None = None) -> dict[str, Any]:
    """Check that the meta-learning pipeline produced expected artifacts."""
    results: dict[str, Any] = {}

    signal_dir = META_DATA_DIR / "signal_buffer"
    signals = list(signal_dir.glob("sig-*.yaml")) if signal_dir.exists() else []
    results["signal_created"] = len(signals) > 0
    results["signal_count"] = len(signals)

    taxonomy_path = META_DATA_DIR / "error_taxonomy.yaml"
    results["taxonomy_updated"] = (
        taxonomy_path.exists() and taxonomy_path.stat().st_size > 50
    )

    skill_path = NANOBOT_SKILLS_DIR / "meta-learning" / "SKILL.md"
    skill_content = ""
    if skill_path.exists():
        try:
            skill_content = skill_path.read_text()
        except OSError:
            pass

    has_rules = _skill_has_learned_rules(skill_content)
    results["skill_updated"] = has_rules
    results["skill_is_bootstrap"] = skill_path.exists() and not has_rules
    results["skill_content_preview"] = skill_content[:500]

    if verify_keywords and skill_content:
        text_lower = skill_content.lower()
        found = [kw for kw in verify_keywords if kw.lower() in text_lower]
        results["keyword_hits"] = found
        results["keyword_match"] = len(found) > 0
    else:
        results["keyword_hits"] = []
        results["keyword_match"] = False

    state_path = META_DATA_DIR / "signal_buffer" / "layer2_state.json"
    if state_path.exists():
        try:
            data = json.loads(state_path.read_text())
            results["layer2_status"] = data.get("status", "unknown")
        except (json.JSONDecodeError, OSError):
            results["layer2_status"] = "read_error"
    else:
        results["layer2_status"] = "not_found"

    results["pipeline_ok"] = all([
        results["signal_created"],
        results["taxonomy_updated"],
        results["skill_updated"],
        results["layer2_status"] == "completed",
    ])
    return results


# =====================================================================
# Cleanup Utilities
# =====================================================================

def _write_bootstrap_skill_md():
    """Write a bootstrap SKILL.md so the agent knows to call capture_signal."""
    from meta_learning.sync_nobot import render_bootstrap_skill_md

    skill_dir = NANOBOT_SKILLS_DIR / "meta-learning"
    skill_dir.mkdir(parents=True, exist_ok=True)
    (skill_dir / "SKILL.md").write_text(render_bootstrap_skill_md(), encoding="utf-8")


def _cleanup_meta_state():
    """Clear signal_buffer, taxonomy, skills, and layer2 state for a fresh scenario."""
    signal_dir = META_DATA_DIR / "signal_buffer"
    if signal_dir.exists():
        for f in signal_dir.glob("sig-*.yaml"):
            f.unlink(missing_ok=True)
        state_file = signal_dir / "layer2_state.json"
        state_file.unlink(missing_ok=True)
        for f in signal_dir.glob("*.processed"):
            f.unlink(missing_ok=True)

    taxonomy_path = META_DATA_DIR / "error_taxonomy.yaml"
    if taxonomy_path.exists():
        taxonomy_path.unlink(missing_ok=True)

    exp_dir = META_DATA_DIR / "experience_pool"
    if exp_dir.exists():
        for f in exp_dir.glob("exp-*.yaml"):
            f.unlink(missing_ok=True)
        idx = exp_dir / "index.yaml"
        idx.unlink(missing_ok=True)

    skills_dir = META_DATA_DIR / "skills"
    if skills_dir.exists():
        shutil.rmtree(skills_dir, ignore_errors=True)
        skills_dir.mkdir(parents=True, exist_ok=True)

    _write_bootstrap_skill_md()

    logger.info("  Meta-learning state cleaned (bootstrap SKILL.md deployed)")


def _extract_response_text(resp: dict[str, Any]) -> str:
    """Best-effort extraction of the assistant's text from a Gateway response."""
    for key in ("content", "response", "message", "text", "answer"):
        if key in resp and isinstance(resp[key], str):
            return resp[key]
    if "choices" in resp:
        try:
            return resp["choices"][0]["message"]["content"]
        except (KeyError, IndexError, TypeError):
            pass
    return json.dumps(resp, ensure_ascii=False)[:4000]


# =====================================================================
# Result Recording
# =====================================================================

def _append_result(result_file: Path, record: dict[str, Any]):
    result_file.parent.mkdir(parents=True, exist_ok=True)
    with open(result_file, "a", encoding="utf-8") as f:
        f.write(json.dumps(record, ensure_ascii=False, default=str) + "\n")


# =====================================================================
# Evaluation Loop
# =====================================================================

async def run_scenario(
    client: DeskClawClient,
    scenario: Scenario,
    result_file: Path,
    api_key: str,
    *,
    skip_teach: bool = False,
    domain_filter: set[str] | None = None,
) -> list[dict[str, Any]]:
    """Run one scenario: cleanup → teach → 5× treatment → return records."""
    records: list[dict[str, Any]] = []
    logger.info("=" * 60)
    logger.info("SCENARIO: %s — %s", scenario.name, scenario.preference)
    logger.info("=" * 60)

    # Phase 0: Setup
    logger.info("--- Phase 0: Setup ---")
    _cleanup_meta_state()
    _prepare_seed_files()

    pipeline_result: dict[str, Any] = {}

    # Phase 1: Teaching
    if not skip_teach:
        logger.info("--- Phase 1: Teaching (%s) ---", scenario.teach_domain)
        session_id = f"teach-{scenario.name}-{uuid.uuid4().hex[:8]}"

        try:
            logger.info("  Sending teach_task …")
            teach_resp = await client.send_message(scenario.teach_task, session_id)
            teach_text = _extract_response_text(teach_resp)
            logger.info("  Teach response: %s", teach_text[:200])

            logger.info("  Sending correction …")
            ts_before = time.time()
            corr_resp = await client.send_message(scenario.correction, session_id)
            corr_text = _extract_response_text(corr_resp)
            logger.info("  Correction response: %s", corr_text[:200])

            logger.info("  Waiting for learning pipeline …")
            pipeline_done = await client.wait_learning_done(
                timeout=180, after_ts=ts_before,
            )
            pipeline_result = check_pipeline_artifacts(scenario.verify_keywords)

            teach_record = {
                "scenario": scenario.name,
                "phase": "teaching",
                "session_id": session_id,
                "teach_response": teach_text[:2000],
                "correction_response": corr_text[:2000],
                "pipeline_done": pipeline_done,
                "pipeline": pipeline_result,
                "timestamp": datetime.now().isoformat(),
            }
            records.append(teach_record)
            _append_result(result_file, teach_record)

            logger.info(
                "  Pipeline: ok=%s  signals=%d  taxonomy=%s  skill=%s  keywords=%s",
                pipeline_result.get("pipeline_ok"),
                pipeline_result.get("signal_count", 0),
                pipeline_result.get("taxonomy_updated"),
                pipeline_result.get("skill_updated"),
                pipeline_result.get("keyword_hits"),
            )
        except Exception as e:
            logger.error("  Teaching failed: %s", e)
            teach_record = {
                "scenario": scenario.name,
                "phase": "teaching",
                "error": str(e),
                "timestamp": datetime.now().isoformat(),
            }
            records.append(teach_record)
            _append_result(result_file, teach_record)
    else:
        logger.info("--- Phase 1: Skipped (--skip-teach) ---")
        pipeline_result = check_pipeline_artifacts(scenario.verify_keywords)

    # Phase 2: Treatment
    logger.info("--- Phase 2: Treatment (%d cases) ---", len(scenario.treatment_cases))
    for case in scenario.treatment_cases:
        if domain_filter and case.domain not in domain_filter:
            logger.info("  [%s] skipped (domain filter)", case.domain)
            continue

        session_id = f"treat-{scenario.name}-{case.domain}-{uuid.uuid4().hex[:8]}"
        logger.info("  [%s] session=%s", case.domain, session_id)
        logger.info("  [%s] prompt: %s", case.domain, case.prompt[:80])

        try:
            t0 = time.time()
            resp = await client.send_message(case.prompt, session_id)
            wall_s = time.time() - t0
            response_text = _extract_response_text(resp)
            logger.info("  [%s] response (%0.1fs): %s", case.domain, wall_s, response_text[:200])

            judge = await llm_judge(scenario.preference, response_text, api_key)
            logger.info(
                "  [%s] judge: compliant=%s  confidence=%.2f  reason=%s",
                case.domain,
                judge["compliant"],
                judge["confidence"],
                judge["reasoning"][:100],
            )

            case_record = {
                "scenario": scenario.name,
                "phase": "treatment",
                "domain": case.domain,
                "session_id": session_id,
                "prompt": case.prompt,
                "response": response_text[:3000],
                "wall_sec": round(wall_s, 1),
                "judge": judge,
                "timestamp": datetime.now().isoformat(),
            }
        except Exception as e:
            logger.error("  [%s] failed: %s", case.domain, e)
            case_record = {
                "scenario": scenario.name,
                "phase": "treatment",
                "domain": case.domain,
                "prompt": case.prompt,
                "error": str(e),
                "judge": {"compliant": None, "confidence": 0.0, "reasoning": f"error: {e}"},
                "timestamp": datetime.now().isoformat(),
            }

        records.append(case_record)
        _append_result(result_file, case_record)

        await client.delete_session(session_id)

    return records


async def run_experiment(args: argparse.Namespace):
    """Main entry: run selected scenarios and produce reports."""
    api_key = os.environ.get("NODESK_API_KEY", "")
    if not api_key:
        logger.error("NODESK_API_KEY not set — LLM Judge will fail")

    run_id = f"deskclaw_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    result_dir = RESULTS_BASE / run_id
    result_dir.mkdir(parents=True, exist_ok=True)
    result_file = result_dir / "results.jsonl"

    client = DeskClawClient(base_url=args.gateway)

    logger.info("Checking DeskClaw Gateway at %s …", args.gateway)
    if not await client.health_check():
        logger.error("Gateway health check failed — is DeskClaw running?")
        await client.close()
        return

    logger.info("Gateway OK. Run ID: %s", run_id)

    scenarios = SCENARIOS
    if args.scenarios:
        names = {n.strip() for n in args.scenarios.split(",")}
        scenarios = [s for s in SCENARIOS if s.name in names]
        if not scenarios:
            logger.error("No matching scenarios: %s", args.scenarios)
            await client.close()
            return
        logger.info("Selected scenarios: %s", [s.name for s in scenarios])

    domain_filter: set[str] | None = None
    if args.domains:
        domain_filter = {d.strip() for d in args.domains.split(",")}
        logger.info("Domain filter: %s", domain_filter)

    all_records: list[dict[str, Any]] = []

    for scenario in scenarios:
        records = await run_scenario(
            client, scenario, result_file, api_key,
            skip_teach=args.skip_teach,
            domain_filter=domain_filter,
        )
        all_records.extend(records)

    await client.close()

    _write_summary(result_dir / "summary.json", all_records)
    logger.info("All results written to %s", result_dir)


# =====================================================================
# Report Generation
# =====================================================================

def _write_summary(path: Path, records: list[dict[str, Any]]):
    """Produce summary.json: scenario × domain matrix + overall stats."""
    scenario_stats: dict[str, dict[str, Any]] = {}

    for scenario in SCENARIOS:
        name = scenario.name
        teach = [r for r in records if r.get("scenario") == name and r.get("phase") == "teaching"]
        treatments = [r for r in records if r.get("scenario") == name and r.get("phase") == "treatment"]

        pipeline_ok = False
        if teach:
            pipeline_ok = teach[0].get("pipeline", {}).get("pipeline_ok", False)

        cases: dict[str, dict[str, Any]] = {}
        for r in treatments:
            domain = r.get("domain", "unknown")
            j = r.get("judge", {})
            cases[domain] = {
                "compliant": j.get("compliant"),
                "confidence": j.get("confidence", 0.0),
                "reasoning": j.get("reasoning", ""),
                "error": r.get("error"),
            }

        compliant_count = sum(1 for c in cases.values() if c["compliant"] is True)
        total_cases = len(cases)

        confidences = [c["confidence"] for c in cases.values() if c["confidence"] is not None]

        scenario_stats[name] = {
            "pipeline_ok": pipeline_ok,
            "cases": cases,
            "compliant_count": compliant_count,
            "total_cases": total_cases,
            "compliance_rate": round(compliant_count / total_cases, 3) if total_cases else 0.0,
            "avg_confidence": round(sum(confidences) / len(confidences), 3) if confidences else 0.0,
        }

    total_cases = sum(s["total_cases"] for s in scenario_stats.values())
    total_compliant = sum(s["compliant_count"] for s in scenario_stats.values())
    pipeline_ok_count = sum(1 for s in scenario_stats.values() if s["pipeline_ok"])

    same_domain_compliant = 0
    same_domain_total = 0
    cross_domain_compliant = 0
    cross_domain_total = 0
    for sc in SCENARIOS:
        stats = scenario_stats.get(sc.name, {})
        for domain, case_data in stats.get("cases", {}).items():
            if case_data.get("compliant") is None:
                continue
            if domain == sc.teach_domain:
                same_domain_total += 1
                if case_data["compliant"]:
                    same_domain_compliant += 1
            else:
                cross_domain_total += 1
                if case_data["compliant"]:
                    cross_domain_compliant += 1

    overall = {
        "total_cases": total_cases,
        "compliant_cases": total_compliant,
        "compliance_rate": round(total_compliant / total_cases, 3) if total_cases else 0.0,
        "pipeline_success_rate": round(pipeline_ok_count / len(SCENARIOS), 3) if SCENARIOS else 0.0,
        "same_domain_transfer_rate": (
            round(same_domain_compliant / same_domain_total, 3) if same_domain_total else 0.0
        ),
        "cross_domain_transfer_rate": (
            round(cross_domain_compliant / cross_domain_total, 3) if cross_domain_total else 0.0
        ),
    }

    summary = {
        "run_timestamp": datetime.now().isoformat(),
        "scenarios": scenario_stats,
        "overall": overall,
    }

    with open(path, "w", encoding="utf-8") as f:
        json.dump(summary, f, indent=2, ensure_ascii=False)

    _print_console_report(scenario_stats, overall)


def _print_console_report(
    scenario_stats: dict[str, dict[str, Any]],
    overall: dict[str, Any],
):
    logger.info("\n" + "=" * 70)
    logger.info("              DESKCLAW META-LEARNING EVALUATION REPORT")
    logger.info("=" * 70)

    header = f"{'Scenario':<22} {'Pipeline':>8} {'Shell':>7} {'Code':>7} {'Excel':>7} {'Word':>7} {'PPT':>7} {'Rate':>7}"
    logger.info(header)
    logger.info("-" * 70)

    for name, stats in scenario_stats.items():
        cases = stats.get("cases", {})

        def _sym(domain: str) -> str:
            c = cases.get(domain, {})
            if c.get("compliant") is True:
                return "PASS"
            elif c.get("compliant") is False:
                return "FAIL"
            return "  - "

        pipe = " OK " if stats["pipeline_ok"] else "FAIL"
        rate = f"{stats['compliance_rate']:.0%}"
        logger.info(
            f"{name:<22} {pipe:>8} {_sym('shell'):>7} {_sym('coding'):>7} "
            f"{_sym('excel'):>7} {_sym('word'):>7} {_sym('pptx'):>7} {rate:>7}"
        )

    logger.info("-" * 70)
    logger.info(
        "Overall: %d/%d compliant (%.0f%%)  |  Pipeline: %.0f%%  |  Cross-domain: %.0f%%",
        overall["compliant_cases"],
        overall["total_cases"],
        overall["compliance_rate"] * 100,
        overall["pipeline_success_rate"] * 100,
        overall["cross_domain_transfer_rate"] * 100,
    )
    logger.info("=" * 70)


# =====================================================================
# CLI
# =====================================================================

def main():
    parser = argparse.ArgumentParser(
        description="DeskClaw Meta-Learning Automated Evaluation",
    )
    parser.add_argument(
        "--gateway", default=DEFAULT_GATEWAY,
        help=f"Gateway URL (default: {DEFAULT_GATEWAY})",
    )
    parser.add_argument(
        "--scenarios", default=None,
        help="Comma-separated scenario names to run (default: all)",
    )
    parser.add_argument(
        "--domains", default=None,
        help="Comma-separated domain names to test (default: all 5)",
    )
    parser.add_argument(
        "--skip-teach", action="store_true",
        help="Skip teaching phase (reuse existing learning state)",
    )
    parser.add_argument(
        "--meta-data-dir", default=None,
        help="Override meta-learning data directory",
    )
    args = parser.parse_args()

    if args.meta_data_dir:
        global META_DATA_DIR
        META_DATA_DIR = Path(args.meta_data_dir).expanduser()

    asyncio.run(run_experiment(args))


if __name__ == "__main__":
    main()
